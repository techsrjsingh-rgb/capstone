"""
Generate a professional PowerPoint presentation for the
Fraud Detection AI Agent capstone project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy – slide background
DARK_CARD  = RGBColor(0x11, 0x27, 0x40)   # card background
ACCENT     = RGBColor(0x00, 0xC2, 0xFF)   # cyan accent
ACCENT2    = RGBColor(0x7C, 0x3A, 0xED)   # purple accent
GREEN      = RGBColor(0x00, 0xE0, 0x96)   # safe / positive
AMBER      = RGBColor(0xFF, 0xB8, 0x00)   # suspicious / warning
RED        = RGBColor(0xFF, 0x3B, 0x30)   # high risk
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xB0, 0xBE, 0xC5)
MID_GREY   = RGBColor(0x55, 0x6B, 0x82)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helper utilities ─────────────────────────────────────────────────────────

def add_slide(prs, layout_index=6):
    layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(layout)


def bg(slide, color=NAVY):
    """Fill slide background with solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill_color=DARK_CARD, line_color=None, line_width=Pt(0)):
    """Add a filled rectangle (card) to the slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def accent_bar(slide, t=0.72, width=13.33, height=0.06, color=ACCENT):
    """Thin horizontal accent bar."""
    rect(slide, 0, t, width, height, fill_color=color)


def txb(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def txb_lines(slide, lines, l, t, w, h,
              size=14, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, line_space=None):
    """Add a text box with multiple lines (list of (text, bold, color, size))."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, b, c, s = item, bold, color, size
        elif len(item) == 2:
            txt, b = item; c, s = color, size
        elif len(item) == 3:
            txt, b, c = item; s = size
        else:
            txt, b, c, s = item

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_space:
            p.space_before = Pt(line_space)
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(s)
        run.font.bold = b
        run.font.color.rgb = c
    return box


def circle(slide, l, t, d, fill_color=ACCENT, text="", text_size=18,
           text_color=NAVY, bold=True):
    """Add a circle with centred text."""
    from pptx.enum.text import PP_ALIGN
    shape = slide.shapes.add_shape(
        9,  # OVAL
        Inches(l), Inches(t), Inches(d), Inches(d)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(text_size)
        run.font.bold = bold
        run.font.color.rgb = text_color
    return shape


def slide_number(slide, n, total=15):
    txb(slide, f"{n} / {total}", 12.5, 7.1, 0.8, 0.3,
        size=10, color=MID_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 15

    # ── 1  TITLE ──────────────────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)

    # gradient-like decorative rectangles
    rect(s, 0, 0, 13.33, 0.5,  fill_color=RGBColor(0x08,0x14,0x22))
    rect(s, 0, 7.0, 13.33, 0.5, fill_color=RGBColor(0x08,0x14,0x22))

    # big cyan side bar
    rect(s, 0, 0, 0.18, 7.5, fill_color=ACCENT)

    # glowing circle decoration (top right)
    circle(s, 10.8, 0.3, 2.2, fill_color=RGBColor(0x05,0x1A,0x35))
    circle(s, 11.1, 0.6, 1.6, fill_color=RGBColor(0x00,0x44,0x77))

    txb(s, "CAPSTONE PROJECT", 0.5, 1.2, 9, 0.5,
        size=13, bold=False, color=ACCENT, italic=True)

    txb(s, "Fraud Detection", 0.5, 1.85, 10, 1.0,
        size=52, bold=True, color=WHITE)
    txb(s, "AI Agent", 0.5, 2.7, 10, 1.0,
        size=52, bold=True, color=ACCENT)

    txb(s, "A Multi-Agent System for Real-Time Banking Fraud Classification",
        0.5, 3.75, 9, 0.5, size=18, color=LIGHT_GREY)

    accent_bar(s, t=4.4, width=5.0, height=0.04, color=ACCENT2)

    txb_lines(s, [
        ("Built with  Anthropic Claude API  ·  Python  ·  Streamlit  ·  FastMCP", False, LIGHT_GREY, 13),
        ("June 2026", False, MID_GREY, 12),
    ], 0.5, 4.6, 9, 0.8, size=13)

    slide_number(s, 1, total)

    # ── 2  AGENDA ─────────────────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)
    accent_bar(s, t=0.78)
    txb(s, "AGENDA", 0.5, 0.18, 5, 0.45, size=11, bold=True, color=ACCENT)
    txb(s, "What We Will Cover Today", 0.5, 0.45, 10, 0.55, size=26, bold=True)

    items = [
        ("01", "Business Problem",          "Why fraud detection needs AI"),
        ("02", "Solution Overview",         "What we built and how it works"),
        ("03", "Architecture",              "4-agent pipeline design"),
        ("04", "Fraud Rules Engine",        "Deterministic Python rules"),
        ("05", "AI Agent Pipeline",         "Sonnet · Opus · Extended Thinking"),
        ("06", "MCP & Tool-Use",            "3 FastMCP servers + tool schemas"),
        ("07", "Hooks & Governance",        "Middleware, audit, compliance"),
        ("08", "Observability",             "Correlation IDs, OTEL, metrics"),
        ("09", "Dashboard Demo",            "Streamlit UI walkthrough"),
        ("10", "Testing & Evaluation",      "79 tests · 100 % pass rate"),
        ("11", "Load Testing",              "Performance & throughput"),
        ("12", "Design Decisions",          "Why we chose this approach"),
        ("13", "Business Impact",           "Quantitative & qualitative value"),
        ("14", "Lessons Learned",           "Key takeaways"),
        ("15", "Thank You & Q&A",           ""),
    ]

    cols = [items[:8], items[8:]]
    xs = [0.45, 6.9]
    for col_idx, col in enumerate(cols):
        x = xs[col_idx]
        for i, (num, title, sub) in enumerate(col):
            y = 1.35 + i * 0.72
            circle(s, x, y + 0.02, 0.42, fill_color=ACCENT2,
                   text=num, text_size=10, text_color=WHITE)
            txb(s, title, x + 0.55, y, 5.8, 0.35, size=14, bold=True)
            if sub:
                txb(s, sub, x + 0.55, y + 0.28, 5.8, 0.3,
                    size=10, color=LIGHT_GREY)

    slide_number(s, 2, total)

    # ── 3  BUSINESS PROBLEM ───────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)
    accent_bar(s)
    txb(s, "BUSINESS PROBLEM", 0.5, 0.18, 5, 0.4, size=11, bold=True, color=ACCENT)
    txb(s, "Why Banking Fraud Needs an AI-Powered Solution", 0.5, 0.45, 12, 0.6,
        size=26, bold=True)

    stats = [
        ("$485B",  "Annual global fraud losses\n(2023 estimate)",  ACCENT),
        ("~1%",    "Transactions reviewed\nmanually",              AMBER),
        ("14 hrs", "Avg. time to detect\nfraud after incident",    RED),
        ("97%",    "False positive rate with\nlegacy rule systems", ACCENT2),
    ]
    for i, (val, label, col) in enumerate(stats):
        x = 0.45 + i * 3.2
        rect(s, x, 1.25, 2.9, 1.8, fill_color=DARK_CARD,
             line_color=col, line_width=Pt(2))
        txb(s, val, x + 0.15, 1.45, 2.6, 0.75,
            size=38, bold=True, color=col, align=PP_ALIGN.CENTER)
        txb(s, label, x + 0.1, 2.18, 2.7, 0.7,
            size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # Gap section
    rect(s, 0.45, 3.3, 12.4, 1.3, fill_color=RGBColor(0x1A,0x0A,0x3A),
         line_color=ACCENT2, line_width=Pt(1.5))
    txb(s, "The Gap", 0.75, 3.42, 4, 0.4, size=14, bold=True, color=ACCENT2)
    txb_lines(s, [
        ("Legacy rule systems  →  Too rigid, high false positives, miss novel patterns", False, LIGHT_GREY, 13),
        ("ML black boxes  →  Opaque, hard to audit, fail regulatory explainability requirements", False, LIGHT_GREY, 13),
        ("Manual review  →  Cannot scale to millions of daily transactions", False, LIGHT_GREY, 13),
    ], 0.75, 3.82, 12.0, 0.9, line_space=4)

    txb(s, "We need a system that is  FAST  +  EXPLAINABLE  +  ACCURATE",
        0.45, 4.75, 12.4, 0.5, size=16, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)

    slide_number(s, 3, total)

    # ── 4  SOLUTION OVERVIEW ──────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)
    accent_bar(s)
    txb(s, "SOLUTION OVERVIEW", 0.5, 0.18, 5, 0.4, size=11, bold=True, color=ACCENT)
    txb(s, "Fraud Detection AI Agent", 0.5, 0.45, 10, 0.55, size=26, bold=True)

    # Left: what it does
    rect(s, 0.4, 1.2, 5.9, 5.6, fill_color=DARK_CARD)
    txb(s, "What It Does", 0.6, 1.35, 5.5, 0.4, size=16, bold=True, color=ACCENT)
    caps = [
        ("< 10 sec",   "Analyzes any transaction"),
        ("0 – 100",    "Numeric risk score"),
        ("3 verdicts", "Safe · Suspicious · High Risk"),
        ("Plain text", "Every decision explained"),
        ("JSONL log",  "All decisions auditable"),
        ("Self-heals", "Retries on API failure"),
    ]
    for i, (kw, desc) in enumerate(caps):
        y = 1.9 + i * 0.75
        rect(s, 0.6, y, 1.2, 0.52, fill_color=ACCENT2)
        txb(s, kw, 0.62, y + 0.06, 1.16, 0.42,
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(s, desc, 1.95, y + 0.1, 4.2, 0.35, size=13, color=LIGHT_GREY)

    # Right: tech stack
    rect(s, 6.8, 1.2, 6.1, 5.6, fill_color=DARK_CARD)
    txb(s, "Tech Stack", 7.0, 1.35, 5.7, 0.4, size=16, bold=True, color=ACCENT)

    stack = [
        ("LLM",            "Anthropic Claude API"),
        ("Models",         "claude-opus-4-8  +  claude-sonnet-4-6"),
        ("Extended Think", "5,000 thinking tokens (RiskScorer)"),
        ("UI",             "Streamlit + Plotly"),
        ("MCP Servers",    "3× FastMCP (ports 8002-8004)"),
        ("Observability",  "OpenTelemetry + Correlation IDs"),
        ("Language",       "Python 3.14"),
        ("Tests",          "pytest + pytest-asyncio  (79 tests)"),
        ("Load Tests",     "Locust"),
    ]
    for i, (layer, tech) in enumerate(stack):
        y = 1.9 + i * 0.54
        txb(s, layer, 7.0, y, 2.1, 0.4, size=11, bold=True, color=ACCENT)
        txb(s, tech, 9.15, y, 3.6, 0.4, size=11, color=LIGHT_GREY)

    slide_number(s, 4, total)

    # ── 5  ARCHITECTURE ───────────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)
    accent_bar(s)
    txb(s, "ARCHITECTURE", 0.5, 0.18, 5, 0.4, size=11, bold=True, color=ACCENT)
    txb(s, "4-Agent Pipeline Architecture", 0.5, 0.45, 12, 0.55, size=26, bold=True)

    # Horizontal pipeline
    stages = [
        ("Streamlit UI",       "frontend/app.py",          RGBColor(0x22,0x2A,0x3F), LIGHT_GREY),
        ("Hooks Middleware",   "pre_process / post_process", RGBColor(0x22,0x2A,0x3F), LIGHT_GREY),
        ("Rules Agent",        "Pure Python · no LLM",      RGBColor(0x00,0x44,0x22), GREEN),
        ("Pattern Agent",      "claude-sonnet-4-6",          RGBColor(0x00,0x2A,0x55), ACCENT),
        ("Risk Scorer",        "claude-opus-4-8\n+ Extended Thinking", RGBColor(0x30,0x11,0x55), ACCENT2),
        ("Coordinator",        "claude-opus-4-8\nFinal Verdict",       RGBColor(0x44,0x11,0x00), RED),
    ]

    bw, bh, gap = 1.9, 1.55, 0.22
    total_w = len(stages) * bw + (len(stages)-1) * gap
    startx = (13.33 - total_w) / 2

    for i, (title, sub, bg_col, col) in enumerate(stages):
        x = startx + i * (bw + gap)
        rect(s, x, 1.4, bw, bh, fill_color=bg_col, line_color=col, line_width=Pt(2))
        txb(s, title, x + 0.08, 1.52, bw - 0.16, 0.45,
            size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
        txb(s, sub, x + 0.08, 2.0, bw - 0.16, 0.85,
            size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
        # Arrow between boxes
        if i < len(stages) - 1:
            ax = x + bw + 0.04
            txb(s, "→", ax, 1.95, gap + 0.1, 0.4, size=16, bold=True,
                color=ACCENT, align=PP_ALIGN.CENTER)

    # Result badges
    txb(s, "Output:", 0.4, 3.3, 1.5, 0.4, size=13, bold=True, color=LIGHT_GREY)
    badges = [("✅  Safe", GREEN), ("⚠️  Suspicious", AMBER), ("🚨  High Risk", RED)]
    for i, (label, col) in enumerate(badges):
        rx = 1.9 + i * 2.5
        rect(s, rx, 3.25, 2.1, 0.5, fill_color=DARK_CARD, line_color=col, line_width=Pt(2))
        txb(s, label, rx + 0.1, 3.32, 1.9, 0.38, size=14, bold=True,
            color=col, align=PP_ALIGN.CENTER)

    # Supporting services
    txb(s, "Supporting Services", 0.4, 4.05, 4, 0.4, size=13, bold=True, color=LIGHT_GREY)
    svcs = [
        ("MCP Server Fraud\nport 8002", ACCENT),
        ("MCP Server Geo\nport 8003", ACCENT),
        ("MCP Orchestrator\nport 8004", ACCENT),
        ("Audit Log\naudit.jsonl", GREEN),
        ("Observability\nOTEL + Metrics", ACCENT2),
        ("RAG Retriever\nTF-IDF", MID_GREY),
    ]
    for i, (label, col) in enumerate(svcs):
        x = 0.4 + i * 2.15
        rect(s, x, 4.5, 1.95, 0.85, fill_color=DARK_CARD, line_color=col, line_width=Pt(1.5))
        txb(s, label, x + 0.08, 4.58, 1.8, 0.72,
            size=10, color=col, align=PP_ALIGN.CENTER)

    slide_number(s, 5, total)

    # ── 6  FRAUD RULES ENGINE ─────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)
    accent_bar(s)
    txb(s, "FRAUD RULES ENGINE", 0.5, 0.18, 6, 0.4, size=11, bold=True, color=ACCENT)
    txb(s, "4 Deterministic Rules — Pure Python, Zero API Cost", 0.5, 0.45, 12, 0.55,
        size=24, bold=True)

    rules = [
        ("01", "High Amount",      "≥ ₹1,00,000",        "35 pts", ACCENT,  "Transaction amount exceeds the high-value threshold set by RBI guidelines"),
        ("02", "Unusual Location", "Fraud hotspot list",  "30 pts", AMBER,   "Location matches known fraud hotspots e.g. Lagos, Panama, Cayman Islands"),
        ("03", "Rapid Succession", "3+ txns in 5 min",   "25 pts", ACCENT2, "Same customer makes ≥3 transactions within a 300-second window"),
        ("04", "International",    "Outside India",       "15 pts", RED,     "Transaction originates from a location outside India"),
    ]

    for i, (num, name, cond, pts, col, desc) in enumerate(rules):
        y = 1.25 + i * 1.42
        rect(s, 0.4, y, 12.5, 1.28, fill_color=DARK_CARD, line_color=col, line_width=Pt(2))
        circle(s, 0.55, y + 0.38, 0.54, fill_color=col, text=num, text_size=14, text_color=NAVY)
        txb(s, name,  1.25, y + 0.08, 3.0, 0.42, size=16, bold=True, color=col)
        txb(s, cond,  1.25, y + 0.52, 3.0, 0.4,  size=13, color=LIGHT_GREY)
        txb(s, desc,  4.4,  y + 0.25, 6.5, 0.65, size=12, color=LIGHT_GREY)
        rect(s, 11.1, y + 0.3, 1.6, 0.5, fill_color=col)
        txb(s, pts,  11.15, y + 0.35, 1.5, 0.4, size=16, bold=True,
            color=NAVY, align=PP_ALIGN.CENTER)

    # Aggregation logic
    rect(s, 0.4, 7.0, 12.5, 0.38, fill_color=RGBColor(0x08,0x14,0x22))
    txb(s, "Aggregation:   0 rules → Safe     1 rule → Suspicious     2+ rules → High Risk",
        0.6, 7.02, 12.0, 0.35, size=13, bold=True, color=WHITE)

    slide_number(s, 6, total)

    # ── 7  AI AGENT PIPELINE ─────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)
    accent_bar(s)
    txb(s, "AI AGENT PIPELINE", 0.5, 0.18, 6, 0.4, size=11, bold=True, color=ACCENT)
    txb(s, "Four Specialist Agents — Each with One Job", 0.5, 0.45, 12, 0.55,
        size=24, bold=True)

    agents = [
        ("Rules Agent",    "core/rules.py",         "Pure Python",    GREEN,
         ["Zero API cost — instant result",
          "4 deterministic fraud rules",
          "Boundary-tested (33 unit tests)",
          "Powers both UI modes"]),
        ("Pattern Agent",  "claude-sonnet-4-6",     "Behavioral AI",  ACCENT,
         ["Interprets transaction context",
          "Asks: is this unusual for this customer?",
          "Matches known fraud patterns",
          "Concise ≤3 sentence analysis"]),
        ("Risk Scorer",    "claude-opus-4-8",        "Extended Think", ACCENT2,
         ["5,000-token thinking budget",
          "Calls calculate_fraud_score tool",
          "Returns precise 0–100 numeric score",
          "Shows internal chain-of-thought"]),
        ("Coordinator",    "claude-opus-4-8",        "Synthesizer",    RED,
         ["Calls generate_fraud_report tool",
          "Synthesizes all 3 agent outputs",
          "Final Safe/Suspicious/High Risk",
          "Recommended action + explanation"]),
    ]

    for i, (name, model, badge, col, bullets) in enumerate(agents):
        x = 0.3 + i * 3.25
        rect(s, x, 1.25, 3.0, 5.75, fill_color=DARK_CARD, line_color=col, line_width=Pt(2))
        # header
        rect(s, x, 1.25, 3.0, 0.6, fill_color=col)
        txb(s, name, x + 0.1, 1.3, 2.8, 0.4,
            size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txb(s, model, x + 0.1, 1.95, 2.8, 0.35,
            size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
        rect(s, x + 0.65, 2.35, 1.7, 0.38, fill_color=RGBColor(0x08,0x14,0x22))
        txb(s, badge, x + 0.65, 2.38, 1.7, 0.3,
            size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
        for j, bullet in enumerate(bullets):
            txb(s, f"• {bullet}", x + 0.1, 2.88 + j * 0.55, 2.8, 0.5,
                size=11, color=LIGHT_GREY)

        # connector arrow
        if i < len(agents) - 1:
            txb(s, "→", x + 3.02, 3.75, 0.22, 0.5,
                size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Self-healing note
    rect(s, 0.3, 7.05, 12.7, 0.35, fill_color=RGBColor(0x08,0x14,0x22))
    txb(s, "Self-Healing:  3 retries with exponential backoff (1s → 2s → 4s)  ·  Fallback from Opus → Sonnet on rate-limit",
        0.5, 7.07, 12.4, 0.3, size=11, color=LIGHT_GREY)

    slide_number(s, 7, total)

    # ── 8  MCP & TOOL-USE ─────────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)
    accent_bar(s)
    txb(s, "MCP & TOOL-USE", 0.5, 0.18, 6, 0.4, size=11, bold=True, color=ACCENT)
    txb(s, "Model Context Protocol + Structured Tool-Use", 0.5, 0.45, 12, 0.55,
        size=24, bold=True)

    # MCP servers left panel
    rect(s, 0.35, 1.2, 6.2, 5.9, fill_color=DARK_CARD)
    txb(s, "3 Independent FastMCP Servers", 0.55, 1.32, 5.8, 0.42,
        size=15, bold=True, color=ACCENT)

    mcps = [
        ("mcp_server_fraud.py", "port 8002", ACCENT, [
            "get_transaction_history(customer_id)",
            "get_fraud_blacklist()",
            "report_fraud_transaction(txn_id)",
            "get_fraud_statistics()",
        ]),
        ("mcp_server_geo.py", "port 8003", GREEN, [
            "get_country_risk_score(location)",
            "check_ip_location(ip)",
            "get_high_risk_regions()",
            "verify_domestic_location(loc)",
        ]),
        ("mcp_server_orchestrator.py", "port 8004", ACCENT2, [
            "Full 4-agent pipeline exposed via MCP",
            "Enables external system integration",
        ]),
    ]

    y = 1.82
    for sname, port, col, tools in mcps:
        rect(s, 0.5, y, 5.9, 0.3, fill_color=col)
        txb(s, f"{sname}   ({port})", 0.6, y + 0.02, 5.7, 0.28,
            size=12, bold=True, color=NAVY)
        for t in tools:
            txb(s, f"    • {t}", 0.55, y + 0.35, 5.8, 0.32, size=11, color=LIGHT_GREY)
            y += 0.32
        y += 0.48

    # Tool-use right panel
    rect(s, 6.85, 1.2, 6.1, 5.9, fill_color=DARK_CARD)
    txb(s, "4 Claude Tool Schemas", 7.05, 1.32, 5.7, 0.42,
        size=15, bold=True, color=ACCENT)

    tools = [
        ("analyze_transaction",   "Trigger all 4 fraud rules",      ACCENT),
        ("check_customer_history","Behavioral context lookup",       GREEN),
        ("calculate_fraud_score", "Return numeric 0–100 score",      ACCENT2),
        ("generate_fraud_report", "Final structured verdict",        RED),
    ]
    for i, (tname, tdesc, col) in enumerate(tools):
        y = 1.85 + i * 1.22
        rect(s, 7.0, y, 5.8, 1.0, fill_color=RGBColor(0x08,0x14,0x22),
             line_color=col, line_width=Pt(1.5))
        txb(s, tname, 7.12, y + 0.06, 5.5, 0.36, size=13, bold=True, color=col)
        txb(s, tdesc, 7.12, y + 0.48, 5.5, 0.38, size=11, color=LIGHT_GREY)

    slide_number(s, 8, total)

    # ── 9  HOOKS & GOVERNANCE ─────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)
    accent_bar(s)
    txb(s, "HOOKS & GOVERNANCE", 0.5, 0.18, 6, 0.4, size=11, bold=True, color=ACCENT)
    txb(s, "Enterprise-Grade Middleware & Compliance Controls", 0.5, 0.45, 12, 0.55,
        size=24, bold=True)

    # Hooks flow
    rect(s, 0.35, 1.2, 5.9, 5.9, fill_color=DARK_CARD)
    txb(s, "FraudHookManager  (core/hooks.py)", 0.55, 1.32, 5.5, 0.4,
        size=14, bold=True, color=ACCENT)

    hooks = [
        ("pre_process()", ACCENT, [
            "✔ Input validation (fields + types)",
            "✔ Rate limiting (token bucket 20/min)",
            "✔ String sanitization & normalization",
            "✔ Correlation ID generation (UUID-v4)",
        ]),
        ("post_process()", GREEN, [
            "✔ Compliance check (valid decision + reason)",
            "✔ Fairness flag (location-only High Risk)",
            "✔ Audit log append (append-only JSONL)",
            "✔ Metrics increment",
        ]),
        ("on_error()", RED, [
            "✔ Returns 'Suspicious' — never crashes",
            "✔ Logs exception to audit trail",
            "✔ Preserves correlation ID in fallback",
        ]),
    ]

    y = 1.85
    for hname, col, items in hooks:
        rect(s, 0.5, y, 5.6, 0.34, fill_color=col)
        txb(s, hname, 0.6, y + 0.02, 5.4, 0.3, size=13, bold=True, color=NAVY)
        for item in items:
            txb(s, item, 0.6, y + 0.42, 5.5, 0.32, size=11, color=LIGHT_GREY)
            y += 0.32
        y += 0.52

    # Governance controls right
    rect(s, 6.55, 1.2, 6.45, 5.9, fill_color=DARK_CARD)
    txb(s, "GovernanceManager  (core/governance.py)", 6.75, 1.32, 6.1, 0.4,
        size=14, bold=True, color=ACCENT)

    govs = [
        ("Input Validation",   ACCENT,  "Required fields · Type checks · Enum validation\nNegative amount rejection · Empty ID rejection"),
        ("Rate Limiting",      AMBER,   "Token bucket algorithm · 20 req / 60 sec\nPrevents API abuse and runaway Streamlit loops"),
        ("Compliance Check",   GREEN,   "Decision must be Safe / Suspicious / High Risk\nReasoning must be ≥ 15 characters"),
        ("Fairness Check",     ACCENT2, "Flags location-only High Risk decisions\nHuman reviewer notified — decision not overridden"),
        ("Audit Logging",      RED,     "Append-only audit.jsonl · One event per line\nTimestamp · Correlation ID · Agent · Action · Data"),
    ]

    y = 1.88
    for gname, col, desc in govs:
        rect(s, 6.7, y, 6.15, 0.84, fill_color=RGBColor(0x08,0x14,0x22),
             line_color=col, line_width=Pt(1.5))
        txb(s, gname, 6.85, y + 0.04, 5.8, 0.28, size=12, bold=True, color=col)
        txb(s, desc, 6.85, y + 0.36, 5.8, 0.45, size=10, color=LIGHT_GREY)
        y += 1.0

    slide_number(s, 9, total)

    # ── 10  OBSERVABILITY ─────────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)
    accent_bar(s)
    txb(s, "OBSERVABILITY", 0.5, 0.18, 6, 0.4, size=11, bold=True, color=ACCENT)
    txb(s, "End-to-End Traceability for Every Transaction", 0.5, 0.45, 12, 0.55,
        size=24, bold=True)

    obs = [
        ("Correlation IDs",     ACCENT,  "UUID-v4 generated per request in pre_process()\nAll 4 agent calls share the same ID\nFull trace retrievable via get_audit_trail(cid)"),
        ("Agent Traces",        ACCENT2, "Agent name · Input summary · Output summary\nDuration in milliseconds · UTC timestamp\nStored in-memory, shown in Streamlit dashboard"),
        ("Metrics",             GREEN,   "fraud_classification  (count per risk level)\nrisk_score  (per transaction, 0–100)\nagent_error  (error rate counter)"),
        ("OpenTelemetry",       AMBER,   "Optional OTEL integration (graceful fallback)\nDistributed tracing spans per agent call\nExported to any OTEL-compatible backend"),
        ("Audit Log  (JSONL)",  RED,     "audit.jsonl  · append-only file\nMachine-readable · streamable · grep-friendly\nRetained indefinitely (rotate externally in prod)"),
        ("Structured Logging",  MID_GREY,"JSON-formatted log lines\nCorrelation ID on every log entry\nLog level: INFO / WARNING / ERROR"),
    ]

    for i, (name, col, desc) in enumerate(obs):
        col_idx = i % 2
        row_idx = i // 2
        x = 0.35 + col_idx * 6.55
        y = 1.25 + row_idx * 2.0
        rect(s, x, y, 6.1, 1.78, fill_color=DARK_CARD, line_color=col, line_width=Pt(2))
        txb(s, name, x + 0.18, y + 0.1, 5.7, 0.38, size=14, bold=True, color=col)
        txb(s, desc, x + 0.18, y + 0.55, 5.7, 1.1, size=11, color=LIGHT_GREY)

    slide_number(s, 10, total)

    # ── 11  DASHBOARD DEMO ────────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)
    accent_bar(s)
    txb(s, "DASHBOARD DEMO", 0.5, 0.18, 6, 0.4, size=11, bold=True, color=ACCENT)
    txb(s, "Streamlit Web Dashboard  (frontend/app.py)", 0.5, 0.45, 12, 0.55,
        size=24, bold=True)

    # Placeholder screenshot boxes
    placeholders = [
        ("Batch Analysis Tab",         "Color-coded results table\n22 sample transactions\nOne-click bulk analysis",  0.35, 1.3),
        ("Risk Distribution Chart",    "Plotly pie chart\nSafe · Suspicious · High Risk\nReal-time refresh",          6.85, 1.3),
        ("Single Transaction Form",    "Custom input form\nPlotly risk gauge\nPer-field validation",                   0.35, 4.2),
        ("Agent Trace Inspector",      "Per-agent output viewer\nExtended thinking visible\nCorrelation ID display",   6.85, 4.2),
    ]

    for title, desc, x, y in placeholders:
        rect(s, x, y, 6.1, 2.65, fill_color=RGBColor(0x0A,0x1A,0x2E),
             line_color=ACCENT, line_width=Pt(1))
        txb(s, f"[ {title} ]", x + 0.15, y + 0.12, 5.8, 0.42,
            size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        txb(s, desc, x + 0.2, y + 0.65, 5.7, 1.8,
            size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
        txb(s, "Run:  streamlit run frontend/app.py --server.port 8501",
            x + 0.2, y + 2.25, 5.7, 0.3,
            size=9, color=MID_GREY, align=PP_ALIGN.CENTER)

    slide_number(s, 11, total)

    # ── 12  TESTING & EVALUATION ─────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)
    accent_bar(s)
    txb(s, "TESTING & EVALUATION", 0.5, 0.18, 6, 0.4, size=11, bold=True, color=ACCENT)
    txb(s, "79 Tests · 3 Layers · 100% Pass Rate", 0.5, 0.45, 12, 0.55,
        size=24, bold=True)

    # 3 test cards
    test_cards = [
        ("test_rules.py",   "33 Tests",  GREEN,
         "Unit tests for all 4 fraud rules",
         ["High Amount (6 boundary cases)",
          "Unusual Location (4 hotspot checks)",
          "Rapid Succession (4 velocity checks)",
          "International (5 location checks)",
          "Aggregation logic (4 tests)",
          "Risk score computation (2 tests)"]),
        ("test_agent.py",   "19 Tests",  ACCENT,
         "Integration tests (mocked Claude)",
         ["Full pipeline analysis",
          "Hook pre/post processing",
          "Tool-use response handling",
          "Error handling & fallback",
          "Retry logic verification",
          "Correlation ID generation"]),
        ("test_shared.py",  "27 Tests",  ACCENT2,
         "Governance & observability",
         ["Input validation (6 cases)",
          "Compliance checks (4 cases)",
          "Fairness detection (2 cases)",
          "Rate limiting (3 cases)",
          "Audit logging (3 cases)",
          "Observability traces (5 cases)"]),
    ]

    for i, (fname, count, col, subtitle, bullets) in enumerate(test_cards):
        x = 0.35 + i * 4.35
        rect(s, x, 1.2, 4.1, 5.95, fill_color=DARK_CARD, line_color=col, line_width=Pt(2))
        rect(s, x, 1.2, 4.1, 0.58, fill_color=col)
        txb(s, fname, x + 0.12, 1.25, 3.85, 0.34,
            size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txb(s, count, x + 1.1, 1.85, 1.9, 0.55,
            size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
        txb(s, subtitle, x + 0.12, 2.48, 3.85, 0.36,
            size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
        for j, bullet in enumerate(bullets):
            txb(s, f"• {bullet}", x + 0.22, 3.0 + j * 0.5, 3.65, 0.45,
                size=11, color=LIGHT_GREY)

    # Accuracy table
    rect(s, 0.35, 7.05, 12.6, 0.35, fill_color=RGBColor(0x08,0x14,0x22))
    txb(s, "Accuracy on 22 sample transactions:   Safe 9/9 (100%)   ·   Suspicious 9/9 (100%)   ·   High Risk 4/4 (100%)",
        0.5, 7.07, 12.3, 0.3, size=12, bold=True, color=WHITE)

    slide_number(s, 12, total)

    # ── 13  LOAD TESTING ─────────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)
    accent_bar(s)
    txb(s, "LOAD TESTING", 0.5, 0.18, 6, 0.4, size=11, bold=True, color=ACCENT)
    txb(s, "Performance & Throughput Results  (Locust)", 0.5, 0.45, 12, 0.55,
        size=24, bold=True)

    # Metrics cards
    metrics = [
        ("~200 txns/sec",   "Rules-Only Throughput",   GREEN),
        ("< 5 ms",          "Rules-Only Avg Latency",  GREEN),
        ("2–5 txns/sec",    "AI Agent Throughput",     ACCENT),
        ("2–8 sec",         "AI Agent Avg Latency",    ACCENT),
        ("0 %",             "Rules Error Rate",        GREEN),
        ("< 1 %",           "AI Agent Error Rate",     AMBER),
    ]
    for i, (val, label, col) in enumerate(metrics):
        col_idx = i % 3
        row_idx = i // 3
        x = 0.4 + col_idx * 4.3
        y = 1.25 + row_idx * 1.55
        rect(s, x, y, 4.0, 1.32, fill_color=DARK_CARD, line_color=col, line_width=Pt(2))
        txb(s, val, x + 0.12, y + 0.1, 3.75, 0.65,
            size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
        txb(s, label, x + 0.12, y + 0.82, 3.75, 0.38,
            size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # Recommendation box
    rect(s, 0.4, 4.5, 12.5, 1.2, fill_color=RGBColor(0x00,0x2A,0x10),
         line_color=GREEN, line_width=Pt(2))
    txb(s, "Production Recommendation", 0.65, 4.6, 6, 0.38, size=14, bold=True, color=GREEN)
    txb_lines(s, [
        ("• Run Rules Engine on 100% of transactions  →  instant triage at 200 txns/sec", False, LIGHT_GREY, 12),
        ("• Escalate Rule-flagged transactions to AI Agent  →  deep analysis only when needed", False, LIGHT_GREY, 12),
        ("• Use async queuing (Celery / Redis) for AI Agent to handle burst load", False, LIGHT_GREY, 12),
    ], 0.65, 5.0, 12.0, 0.85, line_space=3)

    # Memory usage
    rect(s, 0.4, 5.88, 5.8, 0.82, fill_color=DARK_CARD,
         line_color=ACCENT2, line_width=Pt(1.5))
    txb(s, "Memory:  Rules-Only < 50 MB   ·   AI Agent < 200 MB",
        0.6, 6.0, 5.4, 0.55, size=12, color=LIGHT_GREY)

    rect(s, 6.55, 5.88, 6.3, 0.82, fill_color=DARK_CARD,
         line_color=ACCENT2, line_width=Pt(1.5))
    txb(s, "Load tool:  locust -f load_tests/locustfile.py --host=http://localhost:8002",
        6.7, 6.0, 6.0, 0.55, size=11, color=LIGHT_GREY)

    slide_number(s, 13, total)

    # ── 14  DESIGN DECISIONS ─────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)
    accent_bar(s)
    txb(s, "DESIGN DECISIONS", 0.5, 0.18, 6, 0.4, size=11, bold=True, color=ACCENT)
    txb(s, "Key Architectural Choices & Rationale", 0.5, 0.45, 12, 0.55,
        size=24, bold=True)

    decisions = [
        ("Pure Python Rules Agent",    ACCENT,  "No LLM dependency means instant, free, 100% testable without API key. Rules engine forms the reliable backbone; AI is reserved for nuanced edge cases."),
        ("Extended Thinking on Opus",  ACCENT2, "5,000 thinking tokens give the Risk Scorer deep reasoning capability. Visible chain-of-thought improves accuracy on borderline transactions."),
        ("Separate MCP Servers",       GREEN,   "Splitting fraud data (port 8002) and geo data (port 8003) allows independent deployment, scaling, and failure isolation."),
        ("Append-Only Audit Log",      AMBER,   "JSONL format is grep-friendly, streamable, and append-only — satisfying financial audit retention requirements without a database."),
        ("Hooks as Middleware",        RED,     "Pre/post hooks cleanly separate cross-cutting concerns (validation, rate-limiting, compliance) from business logic — OCP principle."),
        ("Correlation IDs",            MID_GREY,"Single UUID ties together all 4 agent calls per transaction, enabling end-to-end trace reconstruction for debugging and compliance."),
    ]

    for i, (title, col, rationale) in enumerate(decisions):
        col_idx = i % 2
        row_idx = i // 2
        x = 0.35 + col_idx * 6.55
        y = 1.25 + row_idx * 2.02
        rect(s, x, y, 6.1, 1.82, fill_color=DARK_CARD, line_color=col, line_width=Pt(2))
        txb(s, title, x + 0.18, y + 0.1, 5.5, 0.38, size=14, bold=True, color=col)
        txb(s, rationale, x + 0.18, y + 0.55, 5.7, 1.15, size=11, color=LIGHT_GREY)

    slide_number(s, 14, total)

    # ── 15  THANK YOU ─────────────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s)

    # Decorations
    rect(s, 0, 0, 0.18, 7.5, fill_color=ACCENT)
    circle(s, 9.5, -0.3, 4.0, fill_color=RGBColor(0x05,0x1A,0x35))
    circle(s, 10.2, 0.4,  2.8, fill_color=RGBColor(0x00,0x44,0x77))
    circle(s, 10.8, 0.9,  1.6, fill_color=RGBColor(0x00,0x77,0xAA))

    txb(s, "Thank You", 0.5, 1.5, 9, 1.05,
        size=58, bold=True, color=WHITE)
    txb(s, "& Questions Welcome", 0.5, 2.55, 9, 0.75,
        size=32, bold=False, color=ACCENT)

    accent_bar(s, t=3.5, width=6.0, height=0.05, color=ACCENT2)

    txb(s, "Fraud Detection AI Agent — Capstone Project", 0.5, 3.7, 9, 0.45,
        size=16, color=LIGHT_GREY)
    txb(s, "Powered by Anthropic Claude API · Python · Streamlit · FastMCP",
        0.5, 4.2, 9, 0.4, size=13, color=MID_GREY)

    # Summary pills
    pills = [
        ("4 Agents",     "Rules · Pattern · Scorer · Coordinator", ACCENT),
        ("79 Tests",     "100 % pass rate",                         GREEN),
        ("3 MCP Servers","Fraud · Geo · Orchestrator",              ACCENT2),
        ("Full Audit",   "Append-only JSONL",                       AMBER),
    ]
    for i, (label, sub, col) in enumerate(pills):
        x = 0.5 + i * 3.0
        rect(s, x, 5.0, 2.7, 1.1, fill_color=DARK_CARD, line_color=col, line_width=Pt(2))
        txb(s, label, x + 0.1, 5.08, 2.5, 0.4, size=15, bold=True, color=col)
        txb(s, sub, x + 0.1, 5.52, 2.5, 0.45, size=11, color=LIGHT_GREY)

    txb(s, "github.com / capstone-project",
        0.5, 6.3, 5, 0.4, size=12, color=MID_GREY)

    slide_number(s, 15, total)

    # ── SAVE ──────────────────────────────────────────────────────────────────
    out = "/home/labuser/Desktop/capstone-project/capstone/docs/fraud_detection_presentation.pptx"
    prs.save(out)
    print(f"Saved → {out}")
    return out


if __name__ == "__main__":
    build()
